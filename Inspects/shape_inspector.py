from pptx import Presentation

def inspect_shapes(pptx_path: str, search_text: str = None, slide_idx: int = None):
    print(f"Inspecting Shapes for: {pptx_path}")
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"Error loading {pptx_path}: {e}")
        return

    EMU_TO_INCH = 1.0 / 914400.0

    slides_to_inspect = []
    if slide_idx is not None:
        if 1 <= slide_idx <= len(prs.slides):
            slides_to_inspect.append((slide_idx, prs.slides[slide_idx - 1]))
        else:
            print(f"Slide index {slide_idx} out of range (1 to {len(prs.slides)}).")
            return
    else:
        for i, slide in enumerate(prs.slides):
            # If search_text is provided, only include slides that contain it
            if search_text:
                found = False
                for shape in slide.shapes:
                    if shape.has_text_frame and search_text in shape.text:
                        found = True
                        break
                if found:
                    slides_to_inspect.append((i + 1, slide))
            else:
                slides_to_inspect.append((i + 1, slide))

    if not slides_to_inspect:
        print("No matching slides found.")
        return

    for s_idx, slide in slides_to_inspect:
        print(f"\n--- Slide {s_idx} ---")
        for j, shape in enumerate(slide.shapes):
            top_in = shape.top * EMU_TO_INCH if shape.top else 0
            left_in = shape.left * EMU_TO_INCH if shape.left else 0
            width_in = shape.width * EMU_TO_INCH if shape.width else 0
            height_in = shape.height * EMU_TO_INCH if shape.height else 0
            
            is_ph = getattr(shape, 'is_placeholder', False)
            ph_info = f" (Placeholder Type: {shape.placeholder_format.type})" if is_ph else ""

            text_preview = "N/A"
            lines = 0
            font_size_pt = "N/A"
            if shape.has_text_frame:
                text = shape.text
                lines = text.count('\n') + text.count('\x0b') + 1 if text else 0
                text_preview = text.replace('\n', '\\n').replace('\x0b', '\\v')[:50]
                
                # Attempt to find first explicit font size
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        if run.font.size:
                            font_size_pt = f"{run.font.size.pt:.1f}pt"
                            break
                    if font_size_pt != "N/A":
                        break

            print(f"  Shape {j}: {shape.name} (type: {shape.shape_type}){ph_info}")
            print(f"    Top: {top_in:.3f} in, Left: {left_in:.3f} in, Width: {width_in:.3f} in, Height: {height_in:.3f} in")
            if shape.has_text_frame:
                print(f"    Text ({lines} lines, Font: {font_size_pt}): '{text_preview}'")
