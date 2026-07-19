from pptx.enum.shapes import MSO_SHAPE
from ..render_context import SlideContext

def render(element, ctx: SlideContext, x, y, w, h) -> float:
    ph = ctx.find_placeholder(getattr(element, 'placeholder', None))
    start_x = ph.left if ph else x
    start_y = ph.top if ph else y
    target_w = ph.width if ph else w
    target_h = ph.height if ph else h

    # 1. Fetch defaults from theme
    node_width = ctx.theme.flow.node_width
    node_height = ctx.theme.flow.node_height
    node_gap = ctx.theme.flow.node_gap
    arrow_min_length = ctx.theme.flow.arrow_min_length
    arrow_thickness = ctx.theme.flow.arrow_thickness
    line_offset = ctx.theme.flow.line_offset

    num_nodes = len(element.nodes)
    node_positions = {}

    def style_flow_node(shape):
        shape.fill.solid()
        shape.fill.fore_color.rgb = ctx.theme.color.flow_fill
        shape.line.color.rgb = ctx.theme.color.flow_line
        for p in shape.text_frame.paragraphs:
            p.font.name = ctx.theme.font.name
            p.font.size = ctx.theme.font.size_body_small
            p.font.color.rgb = ctx.theme.color.flow_text

    def style_flow_arrow(shape):
        shape.fill.solid()
        shape.fill.fore_color.rgb = ctx.theme.color.flow_line
        shape.line.color.rgb = ctx.theme.color.flow_line

    if element.direction == 'horizontal':
        # --- Horizontal Dynamic Scaling ---
        total_req_w = num_nodes * node_width + (num_nodes - 1) * node_gap
        if total_req_w > target_w and num_nodes > 0:
            # Min allowed gap: max of 0.15" or arrow_min_length
            min_node_gap = max(int(914400 * 0.15), arrow_min_length)
            scale = target_w / total_req_w
            
            node_gap = max(min_node_gap, int(node_gap * scale))
            remaining_w = target_w - (num_nodes - 1) * node_gap
            
            # If still exceeding, compress node_gap to absolute minimum (arrow_min_length)
            if remaining_w < num_nodes * int(914400 * 0.5):
                node_gap = max(arrow_min_length, int(914400 * 0.1))
                remaining_w = target_w - (num_nodes - 1) * node_gap
                
            node_width = max(int(914400 * 0.4), int(remaining_w / num_nodes))

        # Render Nodes
        for i, node in enumerate(element.nodes):
            nx = start_x + i * (node_width + node_gap)
            ny = start_y
            shape = ctx.slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                nx, ny, node_width, node_height
            )
            shape.text = node.label
            style_flow_node(shape)
            node_positions[f"{node.id}_out"] = (nx + node_width, ny + node_height / 2)
            node_positions[f"{node.id}_in"] = (nx, ny + node_height / 2)

        # Render Arrows
        for edge in element.edges:
            if f"{edge.from_node}_out" in node_positions and f"{edge.to_node}_in" in node_positions:
                fx, fy = node_positions[f"{edge.from_node}_out"]
                tx, ty = node_positions[f"{edge.to_node}_in"]

                if tx >= fx:
                    arrow_x = fx
                    arrow_width = max(tx - fx, arrow_min_length)
                    shape_type = MSO_SHAPE.RIGHT_ARROW
                else:
                    arrow_x = tx
                    arrow_width = max(fx - tx, arrow_min_length)
                    shape_type = MSO_SHAPE.LEFT_ARROW

                arrow = ctx.slide.shapes.add_shape(
                    shape_type,
                    arrow_x, fy - line_offset, arrow_width, arrow_thickness
                )
                style_flow_arrow(arrow)

    else:
        # --- Vertical Dynamic Scaling ---
        total_req_h = num_nodes * node_height + (num_nodes - 1) * node_gap
        if total_req_h > target_h and num_nodes > 0:
            min_node_gap = max(int(914400 * 0.15), arrow_min_length)
            
            # 1. Compress node_gap first, keeping node_height intact
            remaining_h = target_h - num_nodes * node_height
            proposed_gap = int(remaining_h / (num_nodes - 1)) if num_nodes > 1 else node_gap
            
            if proposed_gap >= min_node_gap:
                node_gap = proposed_gap
            else:
                # 2. If gap hits minimum safety limit, also compress node_height
                node_gap = min_node_gap
                remaining_h = target_h - (num_nodes - 1) * node_gap
                node_height = max(int(914400 * 0.35), int(remaining_h / num_nodes))
                
                # Extreme fallback if still exceeding
                if remaining_h < num_nodes * int(914400 * 0.35):
                    node_gap = max(arrow_min_length, int(914400 * 0.1))
                    remaining_h = target_h - (num_nodes - 1) * node_gap
                    node_height = max(int(914400 * 0.3), int(remaining_h / num_nodes))

        # Render Nodes
        for i, node in enumerate(element.nodes):
            nx = start_x
            ny = start_y + i * (node_height + node_gap)
            shape = ctx.slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                nx, ny, node_width, node_height
            )
            shape.text = node.label
            style_flow_node(shape)
            node_positions[f"{node.id}_out"] = (nx + node_width / 2, ny + node_height)
            node_positions[f"{node.id}_in"] = (nx + node_width / 2, ny)

        # Render Arrows
        v_arrow_thickness = int(arrow_thickness * 1.5)
        for edge in element.edges:
            if f"{edge.from_node}_out" in node_positions and f"{edge.to_node}_in" in node_positions:
                fx, fy = node_positions[f"{edge.from_node}_out"]
                tx, ty = node_positions[f"{edge.to_node}_in"]
                if ty >= fy:
                    raw_height = ty - fy
                    arrow_height = max(int(raw_height * 0.8), arrow_min_length)
                    center_y = (fy + ty) / 2
                    arrow_y = center_y - arrow_height / 2
                    shape_type = MSO_SHAPE.DOWN_ARROW
                else:
                    raw_height = fy - ty
                    arrow_height = max(int(raw_height * 0.8), arrow_min_length)
                    center_y = (fy + ty) / 2
                    arrow_y = center_y - arrow_height / 2
                    shape_type = MSO_SHAPE.UP_ARROW

                arrow = ctx.slide.shapes.add_shape(
                    shape_type,
                    fx - v_arrow_thickness / 2, arrow_y, v_arrow_thickness, arrow_height
                )
                style_flow_arrow(arrow)

    if not ph:
        rendered_height = getattr(element, 'height_hint', None)
        if rendered_height is None:
            if element.direction == 'horizontal':
                rendered_height = node_height
            else:
                rendered_height = num_nodes * node_height + (num_nodes - 1) * node_gap
        return y + rendered_height + ctx.theme.layout.element_gap
    return y
