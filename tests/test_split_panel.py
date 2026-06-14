import pytest
from deck2pptx.models import Deck, Slide, Text, BulletList, Split, Panel
from deck2pptx.markdown_adapter import load_markdown
from deck2pptx.yaml_adapter import load_yaml
from deck2pptx.renderer import render_deck
from pptx import Presentation

@pytest.fixture
def temp_dir(tmp_path):
    return tmp_path

def test_markdown_parser_new_syntax(temp_dir):
    md_content = """---
title: Test Deck
---
# Slide 1
<!-- layout "TitleLayout"; subtitle "My Sub" -->
<!-- placeholder "Body" -->
Some text

# Slide 2
<!-- l "SectionLayout" -->
<!-- sub "Sub 2" -->
<!-- ph "Content" -->
- Bullet 1
"""
    md_file = temp_dir / "test.md"
    md_file.write_text(md_content, encoding='utf-8')
    
    deck = load_markdown(md_file)
    
    assert len(deck.slides) == 2
    
    s1 = deck.slides[0]
    assert s1.title == "Slide 1"
    assert s1.layout_hint == "TitleLayout"
    assert s1.subtitle == "My Sub"
    assert len(s1.elements) == 1
    assert isinstance(s1.elements[0], Text)
    assert s1.elements[0].placeholder == "Body"
    
    s2 = deck.slides[1]
    assert s2.title == "Slide 2"
    assert s2.layout_hint == "SectionLayout"
    assert s2.subtitle == "Sub 2"
    assert len(s2.elements) == 1
    assert isinstance(s2.elements[0], BulletList)
    assert s2.elements[0].placeholder == "Content"

def test_markdown_split_horizontal(temp_dir):
    md_content = """# Split Horizontal
<!-- split h -->
<!-- panel "Left" -->
Left text
<!-- panel "Right" -->
- Right item
<!-- /split -->
"""
    md_file = temp_dir / "test_split_h.md"
    md_file.write_text(md_content, encoding='utf-8')
    
    deck = load_markdown(md_file)
    assert len(deck.slides) == 1
    s = deck.slides[0]
    assert len(s.elements) == 1
    split = s.elements[0]
    assert isinstance(split, Split)
    assert split.direction == "horizontal"
    assert len(split.panels) == 2
    
    p1 = split.panels[0]
    assert p1.title == "Left"
    assert len(p1.elements) == 1
    assert isinstance(p1.elements[0], Text)
    
    p2 = split.panels[1]
    assert p2.title == "Right"
    assert len(p2.elements) == 1
    assert isinstance(p2.elements[0], BulletList)

def test_markdown_split_vertical(temp_dir):
    md_content = """# Split Vertical
<!-- split v -->
<!-- panel -->
Top text
<!-- panel -->
Bottom text
<!-- /split -->
"""
    md_file = temp_dir / "test_split_v.md"
    md_file.write_text(md_content, encoding='utf-8')
    
    deck = load_markdown(md_file)
    s = deck.slides[0]
    split = s.elements[0]
    assert isinstance(split, Split)
    assert split.direction == "vertical"
    assert len(split.panels) == 2
    assert split.panels[0].title is None
    assert split.panels[1].title is None

def test_yaml_split(temp_dir):
    yaml_content = """
slides:
  - title: YAML Split
    elements:
      - split:
          direction: horizontal
          panels:
            - title: Panel 1
              elements:
                - text: T1
            - title: Panel 2
              elements:
                - text: T2
"""
    yaml_file = temp_dir / "test_split.yaml"
    yaml_file.write_text(yaml_content, encoding='utf-8')
    
    deck = load_yaml(yaml_file)
    s = deck.slides[0]
    split = s.elements[0]
    assert isinstance(split, Split)
    assert split.direction == "horizontal"
    assert len(split.panels) == 2
    assert split.panels[0].title == "Panel 1"
    assert isinstance(split.panels[0].elements[0], Text)

def test_validation_rejects_nested_split():
    from deck2pptx.validation import validate_deck, ValidationError
    
    deck = Deck()
    split1 = Split(direction="horizontal", panels=[
        Panel(elements=[
            Split(direction="vertical", panels=[Panel(elements=[Text("T1")])])
        ])
    ])
    deck.slides.append(Slide(title="Nested", elements=[split1]))
    
    import pytest
    with pytest.raises(ValidationError) as exc_info:
        validate_deck(deck, ".")
    
    errors = getattr(exc_info.value, "errors", [])
    assert len(errors) > 0
    assert "Nested Split elements are not supported" in errors[0]["message"]

def test_validation_panel_empty():
    from deck2pptx.validation import validate_deck, ValidationError
    
    deck = Deck()
    split = Split(direction="horizontal", panels=[
        Panel(),  # No title, no elements (invalid)
        Panel(title="Valid Title"),  # Title, no elements (valid)
        Panel(elements=[Text("Valid Content")])  # Elements, no title (valid)
    ])
    deck.slides.append(Slide(title="Empty Panel Test", elements=[split]))
    
    import pytest
    with pytest.raises(ValidationError) as exc_info:
        validate_deck(deck, ".")
    
    errors = getattr(exc_info.value, "errors", [])
    assert len(errors) == 1
    assert "Panel 1 is empty" in errors[0]["message"]

def test_renderer_builds_split_pptx(temp_dir):
    deck = Deck()
    split = Split(direction="horizontal", panels=[
        Panel(title="Left", elements=[Text("L Content")]),
        Panel(title="Right", elements=[BulletList(items=["R Item"])])
    ])
    deck.slides.append(Slide(title="Split Test", elements=[split]))
    
    output_path = temp_dir / "output.pptx"
    render_deck(deck, str(output_path), base_dir=temp_dir)
    
    assert output_path.exists()
    
    prs = Presentation(output_path)
    slide = prs.slides[0]
    
    left_shapes = []
    right_shapes = []
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "Left" in text or "L Content" in text:
                left_shapes.append(shape)
            elif "Right" in text or "R Item" in text:
                right_shapes.append(shape)
                
    assert len(left_shapes) >= 2
    assert len(right_shapes) >= 2
    
    for l_shape in left_shapes:
        for r_shape in right_shapes:
            assert l_shape.left < r_shape.left

def test_renderer_builds_vertical_split_pptx(temp_dir):
    deck = Deck()
    split = Split(direction="vertical", panels=[
        Panel(title="Top", elements=[Text("Top Content")]),
        Panel(title="Bottom", elements=[Text("Bottom Content")])
    ])
    deck.slides.append(Slide(title="Split Vertical Test", elements=[split]))
    
    output_path = temp_dir / "output_v.pptx"
    render_deck(deck, str(output_path), base_dir=temp_dir)
    
    assert output_path.exists()
    
    prs = Presentation(output_path)
    slide = prs.slides[0]
    
    top_shapes = []
    bottom_shapes = []
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "Top" in text or "Top Content" in text:
                top_shapes.append(shape)
            elif "Bottom" in text or "Bottom Content" in text:
                bottom_shapes.append(shape)
                
    assert len(top_shapes) >= 2
    assert len(bottom_shapes) >= 2
    
    for t_shape in top_shapes:
        for b_shape in bottom_shapes:
            assert t_shape.top < b_shape.top
