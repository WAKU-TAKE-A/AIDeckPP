import sys
from pathlib import Path
from pptx import Presentation

pptx_path = Path("Inputs/template_tate_01.pptx")
if not pptx_path.exists():
    print(f"File not found: {pptx_path.absolute()}")
    sys.exit(1)

prs = Presentation(str(pptx_path))
print(f"Inspecting {pptx_path.name}")
slide = prs.slides[0]
print(f"Slide 1 has {len(slide.shapes)} shapes.")

for j, shape in enumerate(slide.shapes):
    print(f"\nShape {j+1}: name='{shape.name}', type={shape.shape_type}")
    if shape.has_text_frame:
        width_inches = shape.width / 914400.0 if shape.width else 0
        height_inches = shape.height / 914400.0 if shape.height else 0
        print(f"  Dimensions: W={width_inches:.2f}\" H={height_inches:.2f}\"")
        text = shape.text
        lines = text.count('\n') + 1 if text else 0
        print(f"  Text Lines: {lines}")
        print(f"  Text Preview: {text.replace(chr(10), ' ')[:50]}")
        
        # Check font sizes
        font_sizes = []
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                if run.font.size:
                    font_sizes.append(run.font.size.pt)
        if font_sizes:
            print(f"  Font sizes (pt): {set(font_sizes)}")
        else:
            print("  Font sizes (pt): Not explicitly set in runs")
    else:
        print("  No text frame")
