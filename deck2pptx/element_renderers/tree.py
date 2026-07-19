from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from ..render_context import SlideContext

def render(element, ctx: SlideContext, x, y, w, h) -> float:
    ph = ctx.find_placeholder(getattr(element, 'placeholder', None))
    start_x = ph.left if ph else x
    start_y = ph.top if ph else y
    target_w = ph.width if ph else w
    target_h = ph.height if ph else h

    # 1. Build flat list to understand max depth and leaf counts
    flat_nodes = []
    def build_flat_list(node, level, parent_idx):
        current_idx = len(flat_nodes)
        flat_nodes.append((node, level, parent_idx))
        for child in node.children:
            build_flat_list(child, level + 1, current_idx)
    build_flat_list(element.root, 0, -1)

    # 2. Fetch theme defaults
    node_width = ctx.theme.tree.node_width
    node_height = ctx.theme.tree.node_height
    vertical_gap = ctx.theme.tree.vertical_gap
    horizontal_gap = ctx.theme.tree.horizontal_gap
    connector_line_width = ctx.theme.tree.connector_line_width

    # 3. Dynamic Scaling (Width & Height)
    from collections import defaultdict
    children_by_parent = defaultdict(list)
    max_level = 0
    for i, (node, level, parent_idx) in enumerate(flat_nodes):
        if parent_idx != -1:
            children_by_parent[parent_idx].append(i)
        if level > max_level:
            max_level = level

    # Estimate Leaf Count
    leaf_indices = [i for i in range(len(flat_nodes)) if i not in children_by_parent]
    leaf_count = len(leaf_indices) if leaf_indices else 1

    # Width Scaling (based on max level)
    total_req_w = max_level * (node_width + horizontal_gap) + node_width
    if total_req_w > target_w and max_level > 0:
        min_horizontal_gap = int(914400 * 0.15) # 0.15 inch min gap
        scale_w = target_w / total_req_w
        
        horizontal_gap = max(min_horizontal_gap, int(horizontal_gap * scale_w))
        remaining_w = target_w - max_level * horizontal_gap
        node_width = max(int(914400 * 0.5), int(remaining_w / (max_level + 1))) # 0.5 inch min node width

    indent_width = node_width + horizontal_gap

    # Height Scaling (based on leaf count)
    total_req_h = leaf_count * node_height + (leaf_count - 1) * vertical_gap
    if total_req_h > target_h and leaf_count > 0:
        min_vertical_gap = int(914400 * 0.1) # 0.1 inch min vertical gap
        
        # Try to compress vertical_gap first, keeping node_height intact
        remaining_h = target_h - leaf_count * node_height
        proposed_gap = int(remaining_h / (leaf_count - 1)) if leaf_count > 1 else vertical_gap
        
        if proposed_gap >= min_vertical_gap:
            vertical_gap = proposed_gap
        else:
            # If gap hits minimum safety limit, also compress node_height
            vertical_gap = min_vertical_gap
            remaining_h = target_h - (leaf_count - 1) * vertical_gap
            node_height = max(int(914400 * 0.25), int(remaining_h / leaf_count)) # 0.25 inch min node height

    # 4. Calculate Y positions
    leaf_counter = 0
    node_y = {}
    def calculate_y(idx):
        nonlocal leaf_counter
        children = children_by_parent[idx]
        if not children:
            ny = start_y + leaf_counter * (node_height + vertical_gap)
            node_y[idx] = ny
            leaf_counter += 1
            return ny
        else:
            child_ys = [calculate_y(c_idx) for c_idx in children]
            ny = sum(child_ys) / len(child_ys)
            node_y[idx] = ny
            return ny
    calculate_y(0)

    # 5. Render Nodes
    node_shapes = {}
    for i, (node, level, parent_idx) in enumerate(flat_nodes):
        nx = start_x + level * indent_width
        ny = node_y[i]

        shape = ctx.slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, nx, ny, node_width, node_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = ctx.theme.color.flow_fill
        shape.line.color.rgb = ctx.theme.color.flow_line
        shape.line.width = connector_line_width

        tf = shape.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        p.text = node.label
        p.font.name = ctx.theme.font.name
        p.font.size = ctx.theme.font.size_body_small
        p.font.color.rgb = ctx.theme.color.flow_text
        p.alignment = PP_ALIGN.CENTER

        node_shapes[i] = shape

    # 6. Render Connectors
    for parent_idx, child_indices in children_by_parent.items():
        for c_idx in child_indices:
            connector = ctx.slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, 0, 0, 0, 0)
            connector.line.color.rgb = ctx.theme.color.flow_line
            connector.line.width = connector_line_width
            connector.begin_connect(node_shapes[parent_idx], 3)
            connector.end_connect(node_shapes[c_idx], 1)

    if not ph:
        rendered_height = getattr(element, 'height_hint', None)
        if rendered_height is None:
            rendered_height = leaf_count * node_height + (leaf_count - 1) * vertical_gap
        return y + rendered_height + ctx.theme.layout.element_gap
    return y
