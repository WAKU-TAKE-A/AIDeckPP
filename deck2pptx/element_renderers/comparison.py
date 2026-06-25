from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from ..render_context import SlideContext
from ..text_utils import count_rendered_lines

def render(element, ctx: SlideContext, x, y, w, h) -> float:
    ph = ctx.find_placeholder(getattr(element, 'placeholder', None))
    num_cols = len(element.columns)
    if num_cols == 0:
        return y

    col_width = (ph.width if ph else w) / num_cols
    start_x = ph.left if ph else x
    start_y = ph.top if ph else y
    title_height = ctx.theme.comparison.title_height

    if element.title:
        tb = ctx.slide.shapes.add_textbox(start_x, start_y, col_width * num_cols, title_height)
        p = tb.text_frame.paragraphs[0]
        p.text = element.title
        p.font.name = ctx.theme.font.name
        p.font.size = ctx.theme.font.size_title
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        start_y += title_height

    calib_h = None
    calib_cpi = None
    if ctx.calibrated_metrics and ctx.theme.font.size_body in ctx.calibrated_metrics:
        calib_h = ctx.calibrated_metrics[ctx.theme.font.size_body]['line_height']
        calib_cpi = ctx.calibrated_metrics[ctx.theme.font.size_body]['chars_per_inch']

    col_width_inches = col_width / ctx.theme.calibration.emu_per_inch
    chars_per_line = max(1, int(col_width_inches * calib_cpi)) if calib_cpi else ctx.theme.comparison.fallback_chars_per_line
    # Shared with Text's height estimate: same "no calibration yet" assumption.
    line_height = calib_h if calib_h else ctx.theme.text.line_height

    max_lines = 0
    for col in element.columns:
        lines = 1
        for item in col.items:
            lines += count_rendered_lines(item, chars_per_line)
        if lines > max_lines:
            max_lines = lines

    box_height = max_lines * line_height + ctx.theme.comparison.padding

    for i, col in enumerate(element.columns):
        col_x = start_x + (i * col_width)

        tb = ctx.slide.shapes.add_textbox(col_x, start_y, col_width, box_height)
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = col.label
        p.font.name = ctx.theme.font.name
        p.font.size = ctx.theme.font.size_body
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        
        for item in col.items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.level = 0
            p.font.name = ctx.theme.font.name
            p.font.size = ctx.theme.font.size_body_semi_small
    if not ph:
        rendered_height = getattr(element, 'height_hint', None)
        if rendered_height is None:
            rendered_height = (title_height if element.title else 0) + box_height
        return y + rendered_height + ctx.theme.layout.element_gap
    return y

def render_timeline(element, ctx: SlideContext, x, y, w, h) -> float:
    ph = ctx.find_placeholder(getattr(element, 'placeholder', None))
    start_x = ph.left if ph else x
    start_y = ph.top if ph else y
    width = ph.width if ph else w

    event_height = ctx.theme.timeline.event_height
    label_width = ctx.theme.timeline.label_width
    line_x_offset = ctx.theme.timeline.line_x_offset
    line_width = ctx.theme.timeline.line_width
    line_v_pad = ctx.theme.timeline.line_vertical_padding
    content_x_offset = ctx.theme.timeline.content_x_offset

    for i, ev in enumerate(element.events):
        ev_y = start_y + (i * event_height)
        # Label
        tb = ctx.slide.shapes.add_textbox(start_x, ev_y, label_width, event_height)
        p = tb.text_frame.paragraphs[0]
        p.text = ev.label
        p.font.name = ctx.theme.font.name
        p.font.size = ctx.theme.font.size_title
        p.font.bold = True
        p.font.color.rgb = ctx.theme.color.flow_line
        p.alignment = PP_ALIGN.RIGHT

        # Line (padded top and bottom by line_v_pad)
        line = ctx.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            start_x + line_x_offset, ev_y + line_v_pad,
            line_width, event_height - (2 * line_v_pad)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = ctx.theme.color.flow_line
        line.line.color.rgb = ctx.theme.color.flow_line

        # Content
        tb = ctx.slide.shapes.add_textbox(start_x + content_x_offset, ev_y, width - content_x_offset, event_height)
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = ev.title
        p.font.name = ctx.theme.font.name
        p.font.size = ctx.theme.font.size_body
        p.font.bold = True
        p.font.color.rgb = ctx.theme.color.flow_text

        if ev.description:
            p2 = tf.add_paragraph()
            p2.text = ev.description
            p2.font.name = ctx.theme.font.name
            p2.font.size = ctx.theme.font.size_body_semi_small
            p2.font.color.rgb = ctx.theme.color.text_light

    if not ph:
        rendered_height = getattr(element, 'height_hint', None)
        if rendered_height is None:
            rendered_height = len(element.events) * event_height
        return y + rendered_height + ctx.theme.layout.element_gap
    return y
