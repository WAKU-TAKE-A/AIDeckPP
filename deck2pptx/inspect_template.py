from pptx import Presentation
from pathlib import Path
import json
import sys

def extract_template_info(pptx_path: str | Path, calib: bool = False) -> dict:
    prs = Presentation(pptx_path)
    layouts_info = []
    
    for idx, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for shape in layout.placeholders:
            ph_format = shape.placeholder_format
            placeholders.append({
                "idx": ph_format.idx,
                "name": shape.name,
                "type": str(ph_format.type)
            })
            
        shapes = []
        for shape in layout.shapes:
            if not shape.is_placeholder:
                shapes.append({
                    "name": shape.name,
                    "type": type(shape).__name__
                })
                
        layouts_info.append({
            "index": idx,
            "name": layout.name,
            "placeholders": placeholders,
            "other_shapes": shapes
        })
        
    result = {
        "template": str(pptx_path),
        "layouts": layouts_info
    }
    
    if calib and len(prs.slides) > 0:
        from .height_estimator import extract_template_metrics
        metrics, title_metrics = extract_template_metrics(prs.slides[0])
        result["calibrated_metrics"] = metrics
        result["title_metrics"] = title_metrics
        
    return result

def inspect_template(pptx_path: str | Path, output_format: str = 'text', calib: bool = False):
    try:
        result = extract_template_info(pptx_path, calib)
    except Exception as e:
        if output_format == 'json':
            print(json.dumps({"error": str(e)}, indent=2))
        else:
            print(f"Error loading template: {e}", file=sys.stderr)
        sys.exit(1)
    
    if output_format == 'json':
        print(json.dumps(result, indent=2))
    else:
        print(f"Template: {result['template']}\n")
        for layout in result['layouts']:
            print(f"Layout {layout['index']}: '{layout['name']}'")
            if layout['placeholders']:
                print("  Placeholders:")
                for ph in layout['placeholders']:
                    print(f"    - [{ph['idx']}] '{ph['name']}' (Type: {ph['type']})")
            if layout['other_shapes']:
                print("  Other Shapes:")
                for sh in layout['other_shapes']:
                    print(f"    - '{sh['name']}' ({sh['type']})")
        if "calibrated_metrics" in result:
            tm = result.get("title_metrics", {})
            print("Title Placeholder Metrics:")
            print(f"  - Font Size: {tm.get('font_size_pt')}pt")
            print(f"  - Top: {tm.get('top_inches'):.2f} inches")
            print(f"  - Left: {tm.get('left_inches'):.2f} inches")
            print(f"  - Width: {tm.get('width_inches'):.2f} inches")
            print(f"  - Height: {tm.get('height_inches'):.2f} inches\n")
            
            print("Calibrated Metrics (from first slide):")
            for fs, data in result["calibrated_metrics"].items():
                print(f"  Font Size {fs}pt:")
                print(f"    - Shape Name: '{data['shape_name']}'")
                print(f"    - Lines: {data['lines']}")
                print(f"    - Chars per line: {data['chars_per_line']}")
                print(f"    - Box Width: {data['box_width_inches']:.2f} inches")
                print(f"    - Box Height: {data['box_height_inches']:.2f} inches")
                print(f"    - CPI: {data['cpi']:.2f}")
                print(f"    - Calculated Line Height: {data['height']} EMU")
            print()
