import sys
from pathlib import Path
from pptx import Presentation

pptx_path = Path("Inputs/template_tate_01.pptx")
prs = Presentation(str(pptx_path))
slide = prs.slides[0]

for j, shape in enumerate(slide.shapes):
    if shape.has_text_frame and j >= 1: # Skip title
        text = shape.text
        lines = text.split('\n')
        print(f"\nShape {j+1} (font size ~= {shape.text_frame.paragraphs[0].runs[0].font.size.pt if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs and shape.text_frame.paragraphs[0].runs[0].font.size else 'Unknown'}pt):")
        for i, line in enumerate(lines):
            print(f"  Line {i+1}: Length = {len(line)} characters | Content: '{line}'")
