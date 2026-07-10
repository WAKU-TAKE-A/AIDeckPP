import pytest
from unittest.mock import MagicMock
from pptx.util import Pt
from deck2pptx.element_renderers.table import render
from deck2pptx.models import Table

def test_render_table_basic():
    # Mock context and theme
    mock_ctx = MagicMock()
    mock_ctx.deck.input_format = "markdown"
    mock_ctx.theme.font.name = "Arial"
    mock_ctx.theme.font.size_body = Pt(18)
    mock_ctx.theme.color.text = "000000"
    mock_ctx.theme.table.header_bg = "CCCCCC"
    mock_ctx.theme.table.header_text = "000000"
    mock_ctx.theme.table.cell_bg = "FFFFFF"
    mock_ctx.theme.table.border_color = "000000"
    mock_ctx.theme.layout.element_gap = Pt(10)
    mock_ctx.level_fonts = {}
    mock_ctx.calibrated_metrics = {18.0: {'height': 25, 'cpi': 2.0}}
    mock_ctx.calibrated_heights = {}
    mock_ctx.theme.calibration.emu_per_inch = 914400
    mock_ctx.theme.calibration.fallback_width_inches = 10
    mock_ctx.find_placeholder.return_value = None
    
    mock_slide = MagicMock()
    mock_ctx.slide = mock_slide
    
    # Mock GraphicFrame (Table)
    mock_shape = MagicMock()
    mock_shape.height = 50
    mock_table = MagicMock()
    mock_shape.table = mock_table
    mock_slide.shapes.add_table.return_value = mock_shape
    
    # Mock rows and cells
    mock_cell_1 = MagicMock()
    mock_cell_2 = MagicMock()
    
    mock_p_1 = MagicMock()
    mock_cell_1.text_frame.paragraphs = [mock_p_1]
    mock_cell_1.text_frame.add_paragraph.return_value = mock_p_1
    
    mock_p_2 = MagicMock()
    mock_cell_2.text_frame.paragraphs = [mock_p_2]
    mock_cell_2.text_frame.add_paragraph.return_value = mock_p_2
    
    mock_table.cell.side_effect = lambda r, c: mock_cell_1 if r == 0 else mock_cell_2
    
    # Element
    table_element = Table(
        headers=["Col1", "Col2"],
        rows=[["Val1", "Val2"]]
    )
    
    # Run render
    y = render(table_element, mock_ctx, x=100, y=100, w=500, h=500)
    
    # Check that add_table was called
    mock_slide.shapes.add_table.assert_called_once()
    assert y > 100 # Should return a new Y coordinate
