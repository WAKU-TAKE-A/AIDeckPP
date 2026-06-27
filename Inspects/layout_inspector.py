from pptx import Presentation

def inspect_layouts(pptx_path: str):
    print(f"Inspecting Layouts for: {pptx_path}")
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"Error loading {pptx_path}: {e}")
        return

    EMU_TO_INCH = 1.0 / 914400.0
    w_in = prs.slide_width * EMU_TO_INCH if prs.slide_width else 0
    h_in = prs.slide_height * EMU_TO_INCH if prs.slide_height else 0

    print(f"Slide Size: Width={prs.slide_width} EMU ({w_in:.3f} in), Height={prs.slide_height} EMU ({h_in:.3f} in)")
    print(f"Total Slide Layouts: {len(prs.slide_layouts)}\n")

    for i, layout in enumerate(prs.slide_layouts):
        print(f"[{i}] Layout: '{layout.name}'")
        if not layout.placeholders:
            print("  (No placeholders)")
        else:
            for shape in layout.placeholders:
                ph_format = shape.placeholder_format
                print(f"  - {shape.name} (type: {ph_format.type})")
        print()
