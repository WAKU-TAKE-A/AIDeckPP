from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from ..render_context import SlideContext
from ..text_utils import count_rendered_lines_weighted

def render(element, ctx: SlideContext, x, y, w, h) -> float:
    ph = ctx.find_placeholder(getattr(element, 'placeholder', None))
    start_x = ph.left if ph else x
    start_y = ph.top if ph else y
    width = ph.width if ph else w

    # Determine font sizes based on level_fonts[0] if available (CodeBlock same rule)
    if ctx.level_fonts and 0 in ctx.level_fonts:
        font_size = Pt(ctx.level_fonts[0])
    else:
        font_size = ctx.theme.font.size_body_small

    col_width_inches = width / 914400.0
    
    # Left border line width
    line_w = ctx.theme.timeline.line_width if hasattr(ctx.theme, 'timeline') else Inches(0.05)
    line_w_inches = line_w / 914400.0
    
    # Left margin is line_w + 0.1", Right margin is 0.1" (total 0.2" + line_w)
    avail_width = col_width_inches - (line_w_inches + 0.2)
    
    fs_pt = font_size.pt if hasattr(font_size, 'pt') else float(font_size)
    
    lines = count_rendered_lines_weighted(element.text if element.text else "", fs_pt, avail_width)
    box_height = Inches(lines * ctx.theme.code.line_height_factor + ctx.theme.code.height_padding)

    # 1. Background Rectangle
    shape = ctx.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, start_x, start_y, width, box_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ctx.theme.color.surface
    shape.line.color.rgb = ctx.theme.color.border

    # 2. Text Frame inside Rectangle (using regular font Calibri/etc.)
    tf = shape.text_frame
    tf.word_wrap = True # Blockquotes should wrap long lines
    p = tf.paragraphs[0]
    p.text = element.text
    p.font.name = ctx.theme.font.name
    p.font.size = font_size
    p.font.color.rgb = ctx.theme.color.text
    p.alignment = PP_ALIGN.LEFT

    # 3. Left Border Line
    # Width equals timeline line_width
    line_h = box_height
    line_x = start_x
    line_y = start_y

    line_shape = ctx.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, line_x, line_y, line_w, line_h)
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = RGBColor(153, 153, 153) # Color: R153 G153 B153
    # Remove border from the vertical line
    line_shape.line.fill.background()

    # Shift text margin so it doesn't overlap the left border
    tf.margin_left = line_w + Inches(0.1)

    if not ph:
        rendered_height = getattr(element, 'height_hint', None)
        if rendered_height is None:
            rendered_height = box_height
        return y + rendered_height + ctx.theme.layout.element_gap
    return y
