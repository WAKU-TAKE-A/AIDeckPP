import pytest
from pathlib import Path
from PIL import Image as PILImage
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.util import Inches
from deck2pptx.adapters import load_deck
from deck2pptx.renderer import render_deck

def write_sample_deck(tmp_path, suffix):
    PILImage.new("RGB", (320, 180), (180, 210, 230)).save(tmp_path / "sample_image.png")
    if suffix == ".yaml":
        sample_path = tmp_path / "sample.deck.yaml"
        sample_path.write_text(
            """
title: Sample Business Presentation
slides:
  - title: Business Strategy 2026
    subtitle: A richer, usable baseline
    layout_hint: title
  - title: Core Principles
    elements:
      - text: "Our approach focuses on the following key areas:"
      - bullet_list:
          - Innovation and Agility
          - Customer Success
  - title: Product Gallery
    elements:
      - gallery:
          images:
            - sample_image.png
            - sample_image.png
            - sample_image.png
  - title: Horizontal Process
    elements:
      - flow:
          direction: horizontal
          nodes:
            - id: n1
              label: Research
            - id: n2
              label: Design
            - id: n3
              label: Build
          edges:
            - from: n1
              to: n2
            - from: n2
              to: n3
  - title: Vertical Drilldown
    elements:
      - flow:
          direction: vertical
          nodes:
            - id: top
              label: Executive Summary
            - id: mid
              label: Detailed Analysis
            - id: bot
              label: Raw Data
          edges:
            - from: top
              to: mid
            - from: mid
              to: bot
""",
            encoding="utf-8",
        )
    else:
        sample_path = tmp_path / "sample.md"
        sample_path.write_text(
            """<!-- layout="title" -->
# Business Strategy 2026
<!-- subtitle="A richer, usable baseline" -->

## Core Principles
Our approach focuses on the following key areas:
- Innovation and Agility
- Customer Success

## Product Gallery
![Image 1](sample_image.png)
![Image 2](sample_image.png)
![Image 3](sample_image.png)

## Horizontal Process
```flow
n1: Research
n2: Design
n3: Build
n1 -> n2
n2 -> n3
```

## Vertical Drilldown
```flow vertical
top: Executive Summary
mid: Detailed Analysis
bot: Raw Data
top -> mid
mid -> bot
```
""",
            encoding="utf-8",
        )
    return sample_path

def run_pptx_readback(tmp_path, sample_path):
    deck = load_deck(sample_path)
    out_path = tmp_path / f'{sample_path.name}.pptx'
    
    # Build
    render_deck(deck, out_path, base_dir=sample_path.parent)
    
    assert out_path.exists()
    
    # Readback
    prs = Presentation(str(out_path))
    assert len(prs.slides) == 5
    
    texts = [shape.text for slide in prs.slides for shape in slide.shapes if getattr(shape, 'has_text_frame', False) and shape.text]
    
    assert any("Business Strategy 2026" in t for t in texts)
    assert any("Our approach focuses on the following key areas:" in t for t in texts)
    
    # Verify Gallery slide (slide index 2)
    gallery_slide = prs.slides[2]
    picture_shapes = [shape for shape in gallery_slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(picture_shapes) == 3
    
    # Check horizontal flow
    flow_slide = prs.slides[3]
    auto_shapes = [shape for shape in flow_slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    rectangles = [shape for shape in auto_shapes if shape.auto_shape_type == MSO_SHAPE.RECTANGLE]
    assert len(rectangles) == 3
    
    # Check vertical flow
    vflow_slide = prs.slides[4]
    auto_shapes_v = [shape for shape in vflow_slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    rectangles_v = [shape for shape in auto_shapes_v if shape.auto_shape_type == MSO_SHAPE.RECTANGLE]
    assert len(rectangles_v) == 3

def test_pptx_readback_yaml(tmp_path):
    run_pptx_readback(tmp_path, write_sample_deck(tmp_path, ".yaml"))

def test_pptx_readback_markdown(tmp_path):
    run_pptx_readback(tmp_path, write_sample_deck(tmp_path, ".md"))
        
def test_reordered_flow_readback(tmp_path):
    sample_path = Path('tests/fixtures/reordered_flow.deck.yaml')
    if not sample_path.exists():
        pytest.skip("reordered_flow.deck.yaml not found")
        
    deck = load_deck(sample_path)
    out_path = tmp_path / 'reordered.pptx'
    
    render_deck(deck, out_path, base_dir=sample_path.parent)
    
    prs = Presentation(str(out_path))
    slide = prs.slides[0]
    
    auto_shapes = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    arrows = [shape for shape in auto_shapes if shape.auto_shape_type in (
        MSO_SHAPE.RIGHT_ARROW, MSO_SHAPE.DOWN_ARROW, MSO_SHAPE.LEFT_ARROW, MSO_SHAPE.UP_ARROW
    )]
    
    # 3 edges defined in reordered_flow.deck.yaml
    assert len(arrows) == 3
