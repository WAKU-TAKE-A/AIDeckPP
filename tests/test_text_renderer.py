import pytest
from unittest.mock import Mock, MagicMock
from pptx.util import Pt
from deck2pptx.render_utils import _set_text_frame_text
from deck2pptx.element_renderers.text import render_bullet
from deck2pptx.models import BulletList, ListItem

def test_set_text_frame_text_inline_formatting():
    # Mock text frame and paragraphs
    mock_tf = MagicMock()
    mock_p = MagicMock()
    mock_tf.paragraphs = [mock_p]
    mock_tf.add_paragraph.return_value = mock_p
    
    # Mock runs
    run_1 = MagicMock()
    run_2 = MagicMock()
    run_3 = MagicMock()
    mock_p.add_run.side_effect = [run_1, run_2, run_3]

    # Test Markdown formatting: "Hello **world**!"
    _set_text_frame_text(mock_tf, "Hello **world**!", input_format="markdown")
    
    assert mock_p.add_run.call_count == 3
    assert run_1.text == "Hello "
    assert run_1.font.bold is not True # Should not be explicitly true if not set, or we can just check it's not True
    assert run_2.text == "world"
    assert run_2.font.bold is True
    assert run_3.text == "!"
    
def test_render_bullet_fallback():
    # Mock context and theme
    mock_ctx = MagicMock()
    mock_ctx.deck.input_format = "markdown"
    mock_ctx.theme.font.name = "Arial"
    mock_ctx.theme.font.size_body = Pt(18)
    mock_ctx.theme.font.size_body_small = Pt(16)
    mock_ctx.theme.color.text = "000000"
    mock_ctx.theme.bullet.bullet_chars = ["-", "*"]
    mock_ctx.theme.bullet.fallback_indent_spaces = 4
    mock_ctx.theme.bullet.indent_per_level = 0.5
    mock_ctx.level_fonts = {}
    mock_ctx.calibrated_metrics = {18.0: {'height': 25, 'cpi': 2.0}, 16.0: {'height': 22, 'cpi': 2.2}}
    mock_ctx.calibrated_heights = {}
    mock_ctx.theme.calibration.emu_per_inch = 914400
    mock_ctx.theme.calibration.fallback_width_inches = 10
    mock_ctx.find_placeholder.return_value = None # Force fallback
    
    mock_slide = MagicMock()
    mock_ctx.slide = mock_slide
    
    mock_txBox = MagicMock()
    mock_slide.shapes.add_textbox.return_value = mock_txBox
    
    mock_tf = MagicMock()
    mock_txBox.text_frame = mock_tf
    mock_p = MagicMock()
    mock_tf.paragraphs = [mock_p]
    mock_tf.add_paragraph.return_value = mock_p
    
    # Run 1 is prefix, Run 2 is text
    run_prefix = MagicMock()
    run_text = MagicMock()
    mock_p.add_run.side_effect = [run_prefix, run_text]

    # Element
    bullet_list = BulletList(items=[ListItem(text="Test", level=0)])
    
    render_bullet(bullet_list, mock_ctx, 0, 0, 100, 100)
    
    assert run_prefix.text == "- "
    assert run_text.text == "Test"
    assert mock_p.font.name == "Arial"
