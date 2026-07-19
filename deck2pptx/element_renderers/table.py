import math
from pptx.util import Inches, Pt
from pptx.oxml.xmlchemy import OxmlElement
from ..render_context import SlideContext
from ..text_utils import count_rendered_lines_weighted

def _set_cell_border(cell, color_hex: str, width_pt: float):
    """Low-level XML helper to apply border style to a cell's XML element."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    
    # 1pt = 12700 EMU
    width_emu = str(int(width_pt * 12700))
    
    # lnL: Left, lnR: Right, lnT: Top, lnB: Bottom
    for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
        # Remove any existing border elements of this type
        for child in list(tcPr):
            if child.tag.endswith(border_name):
                tcPr.remove(child)
        
        # Create a new border element
        ln = OxmlElement(f'a:{border_name}')
        ln.set('w', width_emu)
        ln.set('cmpd', 's')  # single line
        
        solidFill = OxmlElement('a:solidFill')
        srgbClr = OxmlElement('a:srgbClr')
        srgbClr.set('val', color_hex)
        solidFill.append(srgbClr)
        ln.append(solidFill)
        
        tcPr.append(ln)

def _get_calibrated_metrics_or_estimate(fs, calibrated_metrics):
    # Check if inputs are mocks or invalid types
    if not isinstance(fs, (int, float)):
        fs = 14.0
        
    if not calibrated_metrics:
        # Estimate based on font size (1 pt = 12700 EMU, line height = 1.2 * font size)
        # CPI estimate: 72.0 / font size (approx 1 char width = font size pt)
        height = int(fs * 12700 * 1.2)
        cpi = 72.0 / fs
        return height, cpi
    if fs in calibrated_metrics:
        return calibrated_metrics[fs]['height'], calibrated_metrics[fs]['cpi']

    # Find closest
    closest_fs = min(calibrated_metrics.keys(), key=lambda k: abs(k - fs))
    ref = calibrated_metrics[closest_fs]
    ratio = fs / closest_fs
    return int(ref['height'] * ratio), ref['cpi'] * (1.0 / ratio)


def render(element, ctx: SlideContext, x, y, w, h) -> float:
    rows = len(element.rows) + 1 if element.headers else len(element.rows)
    cols = len(element.headers) if element.headers else (len(element.rows[0]) if element.rows else 1)

    ph = ctx.find_placeholder(getattr(element, 'placeholder', None))
    target_x = ph.left if ph else x
    target_y = ph.top if ph else y
    target_w = ph.width if ph else w

    table_shape = ctx.slide.shapes.add_table(rows, cols, target_x, target_y, target_w, ctx.theme.table.placeholder_init_height)
    table = table_shape.table

    # Clear template default style and apply Table Grid style
    try:
        table.style = 'Table Grid'
    except Exception:
        pass

    col_width_inches = (target_w / 914400.0) / cols
    margin_top_emu = int(Inches(0.05))
    margin_bottom_emu = int(Inches(0.05))
    margin_left_emu = int(Inches(0.1))
    margin_right_emu = int(Inches(0.1))

    print(f"[DEBUG] Table Height Calculation (cols={cols}, col_width={col_width_inches:.3f} in):")

    row_offset = 0
    if element.headers:
        header_fs = ctx.theme.table.header_font_size.pt if hasattr(ctx.theme.table.header_font_size, 'pt') else ctx.theme.table.header_font_size
        h_line_height, _ = _get_calibrated_metrics_or_estimate(header_fs, ctx.calibrated_metrics)
        
        max_lines = 1
        for i, header in enumerate(element.headers):
            cell_obj = table.cell(0, i)
            cell_obj.text = header
            cell_obj.fill.solid()
            cell_obj.fill.fore_color.rgb = ctx.theme.table.header_fill_color
            
            # Apply custom border
            _set_cell_border(cell_obj, ctx.theme.table.border_color, ctx.theme.table.border_width_pt)
            
            # Explicitly set cell margins
            cell_obj.margin_top = margin_top_emu
            cell_obj.margin_bottom = margin_bottom_emu
            cell_obj.margin_left = margin_left_emu
            cell_obj.margin_right = margin_right_emu

            lines = count_rendered_lines_weighted(str(header), header_fs, col_width_inches)
            if lines > max_lines:
                max_lines = lines
                
            for p in cell_obj.text_frame.paragraphs:
                p.font.name = ctx.theme.font.name
                p.font.size = ctx.theme.table.header_font_size
                p.font.bold = True
                p.font.color.rgb = ctx.theme.color.text
                
        border_emu = int(ctx.theme.table.border_width_pt * 12700)
        row_height = int(max_lines * h_line_height + margin_top_emu + margin_bottom_emu + border_emu)
        table.rows[0].height = row_height
        
        print(f"  - Row 0 (Header): max_lines={max_lines}, fs={header_fs}pt, line_height={h_line_height/914400.0:.3f} in, margins={(margin_top_emu+margin_bottom_emu)/914400.0:.3f} in => row_height={row_height/914400.0:.3f} in")
        row_offset = 1

    cell_fs = ctx.theme.table.cell_font_size.pt if hasattr(ctx.theme.table.cell_font_size, 'pt') else ctx.theme.table.cell_font_size
    c_line_height, _ = _get_calibrated_metrics_or_estimate(cell_fs, ctx.calibrated_metrics)

    for r_idx, row in enumerate(element.rows):
        max_lines = 1
        for c_idx, cell in enumerate(row):
            cell_obj = table.cell(r_idx + row_offset, c_idx)
            cell_obj.text = str(cell)
            cell_obj.fill.solid()
            cell_obj.fill.fore_color.rgb = ctx.theme.color.background
            
            # Apply custom border
            _set_cell_border(cell_obj, ctx.theme.table.border_color, ctx.theme.table.border_width_pt)
            
            # Explicitly set cell margins
            cell_obj.margin_top = margin_top_emu
            cell_obj.margin_bottom = margin_bottom_emu
            cell_obj.margin_left = margin_left_emu
            cell_obj.margin_right = margin_right_emu

            lines = count_rendered_lines_weighted(str(cell), cell_fs, col_width_inches)
            if lines > max_lines:
                max_lines = lines
                
            for p in cell_obj.text_frame.paragraphs:
                p.font.name = ctx.theme.font.name
                p.font.size = ctx.theme.table.cell_font_size
                p.font.color.rgb = ctx.theme.color.text
                
        border_emu = int(ctx.theme.table.border_width_pt * 12700)
        row_height = int(max_lines * c_line_height + margin_top_emu + margin_bottom_emu + border_emu)
        table.rows[r_idx + row_offset].height = row_height
        
        print(f"  - Row {r_idx + row_offset} (Data): max_lines={max_lines}, fs={cell_fs}pt, line_height={c_line_height/914400.0:.3f} in, margins={(margin_top_emu+margin_bottom_emu)/914400.0:.3f} in => row_height={row_height/914400.0:.3f} in")

    total_calculated_height = sum(r.height for r in table.rows)
    print(f"  => Total calculated height: {total_calculated_height/914400.0:.3f} in (EMU: {total_calculated_height})")

    if not ph:
        rendered_height = getattr(element, 'height_hint', None)
        if rendered_height is None:
            rendered_height = table_shape.height
        return y + rendered_height + ctx.theme.layout.element_gap
    return y
