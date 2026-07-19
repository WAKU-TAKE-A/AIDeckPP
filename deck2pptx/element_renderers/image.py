from pptx.enum.text import PP_ALIGN
from ..render_context import SlideContext

def render(element, ctx: SlideContext, x, y, w, h) -> float:
    img_path = ctx.base_dir / element.source
    ph = ctx.find_placeholder(getattr(element, 'placeholder', None))
    caption_height = ctx.theme.image.caption_height
    try:
        if ph and hasattr(ph, 'insert_picture'):
            pic = ph.insert_picture(str(img_path))
            if element.caption:
                tb = ctx.slide.shapes.add_textbox(pic.left, pic.top + pic.height, pic.width, caption_height)
                p = tb.text_frame.paragraphs[0]
                p.text = element.caption
                p.font.name = ctx.theme.font.name
                p.font.size = ctx.theme.font.size_body_extra_small
                p.alignment = PP_ALIGN.CENTER
                p.font.color.rgb = ctx.theme.color.text_light
            return y
        else:
            target_x = ph.left if ph else x
            target_y = ph.top if ph else y
            max_w = ph.width if ph else w
            max_h = ph.height if ph else h
            if element.caption:
                max_h -= caption_height

            # Constrain to minimum safety dimensions
            max_h = max(int(914400 * 0.5), max_h)
            max_w = max(int(914400 * 0.5), max_w)

            # Pre-calculate precise scale using Pillow to avoid python-pptx fallback bugs
            try:
                from PIL import Image as PILImage
                with PILImage.open(img_path) as pil_img:
                    w_px, h_px = pil_img.size
                native_w = w_px * 9525
                native_h = h_px * 9525
                scale = min(max_w / native_w, max_h / native_h, 1.0)
                new_w = int(native_w * scale)
                new_h = int(native_h * scale)
            except Exception:
                new_w = max_w
                new_h = max_h

            pic = ctx.slide.shapes.add_picture(str(img_path), target_x, target_y, width=new_w, height=new_h)

            if element.caption:
                tb = ctx.slide.shapes.add_textbox(target_x, target_y + new_h, new_w, caption_height)
                p = tb.text_frame.paragraphs[0]
                p.text = element.caption
                p.font.name = ctx.theme.font.name
                p.font.size = ctx.theme.font.size_body_extra_small
                p.alignment = PP_ALIGN.CENTER
                p.font.color.rgb = ctx.theme.color.text_light

            if not ph:
                rendered_height = getattr(element, 'height_hint', None)
                if rendered_height is None:
                    rendered_height = new_h + (caption_height if element.caption else 0)
                return y + rendered_height + ctx.theme.layout.element_gap
            return y
    except Exception as e:
        print(f"Failed to load image {img_path}: {e}")
        return y

def render_gallery(element, ctx: SlideContext, x, y, w, h) -> float:
    num_images = len(element.images)
    if num_images == 0:
        return y

    ph = ctx.find_placeholder(getattr(element, 'placeholder', None))
    cols = getattr(element, 'columns', None)
    rows_ct = getattr(element, 'rows', None)
    caption_height = ctx.theme.image.caption_height
    padding = ctx.theme.image.gallery_padding

    if cols is None and rows_ct is None:
        if num_images == 1: cols, rows_ct = 1, 1
        elif num_images == 2: cols, rows_ct = 2, 1
        elif num_images == 3: cols, rows_ct = 3, 1
        elif num_images == 4: cols, rows_ct = 2, 2
        else: cols, rows_ct = 3, 2
    elif cols is None:
        import math
        cols = math.ceil(num_images / rows_ct)
    elif rows_ct is None:
        import math
        rows_ct = math.ceil(num_images / cols)

    cell_width = (ph.width if ph else w) / cols
    cell_height = (ph.height if ph else h) / rows_ct
    start_y = ph.top if ph else y

    max_bottom_y = start_y

    for i, img in enumerate(element.images[:cols*rows_ct]):
        r = i // cols
        c = i % cols
        cell_x = (ph.left if ph else x) + (c * cell_width)
        cell_y = (ph.top if ph else y) + (r * cell_height)

        max_w = cell_width - padding
        max_h = cell_height - padding
        if img.caption:
            max_h -= caption_height

        img_path = ctx.base_dir / img.source
        try:
            pic = ctx.slide.shapes.add_picture(str(img_path), 0, 0)
            native_w = pic.width
            native_h = pic.height
            scale = min(max_w / native_w, max_h / native_h, 1.0)
            new_w = int(native_w * scale)
            new_h = int(native_h * scale)

            # Align top-left instead of centering within the cell (preserving padding)
            left_x = int(cell_x + padding / 2)
            top_y = int(cell_y + padding / 2)

            pic.left = left_x
            pic.top = top_y
            pic.width = new_w
            pic.height = new_h

            bottom_y = top_y + new_h
            if img.caption:
                tb = ctx.slide.shapes.add_textbox(left_x, top_y + new_h, new_w, caption_height)
                p = tb.text_frame.paragraphs[0]
                p.text = img.caption
                p.font.name = ctx.theme.font.name
                p.font.size = ctx.theme.font.size_body_extra_small
                p.alignment = PP_ALIGN.CENTER
                p.font.color.rgb = ctx.theme.color.text_light
                bottom_y += caption_height
                
            if bottom_y > max_bottom_y:
                max_bottom_y = bottom_y
        except Exception as e:
            print(f"Failed to load image {img_path}: {e}")

    if not ph:
        rendered_height = getattr(element, 'height_hint', None)
        if rendered_height is None:
            rendered_height = max_bottom_y - start_y
        return y + rendered_height + ctx.theme.layout.element_gap
    return y
