import pytest
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from deck2pptx.yaml_adapter import load_yaml
from deck2pptx.markdown_adapter import load_markdown
from deck2pptx.models import Comparison, Timeline, CodeBlock, Tree
from deck2pptx.adapters import load_deck
from deck2pptx.renderer import render_deck

def test_yaml_business_elements(tmp_path):
    content = """
title: Test
slides:
  - title: Slide 1
    elements:
      - comparison:
          columns:
            - label: Left
              items: ["a", "b"]
            - label: Right
              items: ["c", "d"]
      - timeline:
          events:
            - label: "2026"
              title: "Phase 1"
              description: "Desc"
      - code_block:
          code: "print('hello')"
          language: "python"
      - tree:
          root:
            label: "A"
            children:
              - label: "B"
"""
    f = tmp_path / "test.yaml"
    f.write_text(content, encoding='utf-8')
    deck = load_yaml(f)
    
    assert len(deck.slides) == 1
    elements = deck.slides[0].elements
    assert len(elements) == 4
    
    comp = elements[0]
    assert isinstance(comp, Comparison)
    assert len(comp.columns) == 2
    assert comp.columns[0].label == "Left"
    assert comp.columns[0].items == ["a", "b"]
    
    tl = elements[1]
    assert isinstance(tl, Timeline)
    assert len(tl.events) == 1
    assert tl.events[0].label == "2026"
    assert tl.events[0].title == "Phase 1"
    
    cb = elements[2]
    assert isinstance(cb, CodeBlock)
    assert cb.code == "print('hello')"
    assert cb.language == "python"
    
    tree = elements[3]
    assert isinstance(tree, Tree)
    assert tree.root.label == "A"
    assert len(tree.root.children) == 1
    assert tree.root.children[0].label == "B"

def test_markdown_business_elements(tmp_path):
    content = """
# Slide 1

```comparison
Left:
- a
- b
Right:
- c
- d
```

```timeline
2026: Phase 1 - Desc
```

```code python
print('hello')
```

```tree
A
  B
```
"""
    f = tmp_path / "test.md"
    f.write_text(content, encoding='utf-8')
    deck = load_markdown(f)
    
    assert len(deck.slides) == 1
    elements = deck.slides[0].elements
    assert len(elements) == 4
    
    comp = elements[0]
    assert isinstance(comp, Comparison)
    assert len(comp.columns) == 2
    assert comp.columns[0].label == "Left"
    assert comp.columns[0].items == ["a", "b"]
    
    tl = elements[1]
    assert isinstance(tl, Timeline)
    assert len(tl.events) == 1
    assert tl.events[0].label == "2026"
    assert tl.events[0].title == "Phase 1"
    assert tl.events[0].description == "Desc"
    
    cb = elements[2]
    assert isinstance(cb, CodeBlock)
    assert cb.code == "print('hello')"
    assert cb.language == "python"
    
    tree = elements[3]
    assert isinstance(tree, Tree)
    assert tree.root.label == "A"
    assert len(tree.root.children) == 1
    assert tree.root.children[0].label == "B"

def test_markdown_tree_uses_block_indent_unit(tmp_path):
    content = """---
indent: 4
---
# Slide 1

```tree
Root
  Child
    Grandchild
```
"""
    f = tmp_path / "tree.md"
    f.write_text(content, encoding='utf-8')
    deck = load_markdown(f)

    tree = deck.slides[0].elements[0]
    assert isinstance(tree, Tree)
    assert tree.root.label == "Root"
    assert tree.root.children[0].label == "Child"
    assert tree.root.children[0].children[0].label == "Grandchild"

def test_business_element_examples_build(tmp_path):
    md_sample = tmp_path / "sample_advanced.md"
    md_sample.write_text(
        """# Advanced

```comparison
Left:
- a
Right:
- b
```

```timeline
2026: Phase 1 - Desc
```

```code python
print("hello")
```

```tree
Root
  Child
```
""",
        encoding="utf-8",
    )
    yaml_sample = tmp_path / "sample_advanced.deck.yaml"
    yaml_sample.write_text(
        """
title: Advanced
slides:
  - title: Advanced
    elements:
      - comparison:
          columns:
            - label: Left
              items: [a]
            - label: Right
              items: [b]
      - timeline:
          events:
            - label: "2026"
              title: "Phase 1"
              description: "Desc"
      - code_block:
          code: print("hello")
          language: python
      - tree:
          root:
            label: Root
            children:
              - label: Child
""",
        encoding="utf-8",
    )

    for sample in [md_sample, yaml_sample]:
        deck = load_deck(sample)
        output = tmp_path / f"{sample.stem}.pptx"
        render_deck(deck, str(output), base_dir=sample.parent)
        assert output.exists()

def test_timeline_uses_blue_monochrome_style(tmp_path):
    content = """
title: Timeline Style
slides:
  - title: Timeline
    elements:
      - timeline:
          events:
            - label: "2026"
              title: "Phase 1"
              description: "Desc"
"""
    sample = tmp_path / "timeline.deck.yaml"
    sample.write_text(content, encoding="utf-8")
    deck = load_deck(sample)
    output = tmp_path / "timeline.pptx"

    render_deck(deck, str(output), base_dir=sample.parent)

    prs = Presentation(str(output))
    slide = prs.slides[0]
    label_shapes = [
        shape for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text == "2026"
    ]
    title_shapes = [
        shape for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and "Phase 1" in shape.text
    ]
    line_shapes = [
        shape for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and shape.auto_shape_type == MSO_SHAPE.RECTANGLE
        and shape.width < shape.height
    ]

    assert label_shapes[0].text_frame.paragraphs[0].font.color.rgb == RGBColor(0x25, 0x63, 0xEB)
    assert title_shapes[0].text_frame.paragraphs[0].font.color.rgb == RGBColor(0x1E, 0x3A, 0x8A)
    assert line_shapes[0].fill.fore_color.rgb == RGBColor(0x25, 0x63, 0xEB)
