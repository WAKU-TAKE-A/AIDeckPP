import sys
import subprocess
import shutil
from pathlib import Path
from pptx import Presentation

def inspect_compare(pptx_path: str, slide_idx: int = 1):
    input_pptx = Path(pptx_path)
    if not input_pptx.exists():
        print(f"Error: File {pptx_path} does not exist.")
        return

    out_dir = Path("Outputs/calibrated")
    out_dir.mkdir(parents=True, exist_ok=True)

    # Run LibreOffice to get the actual recalculated heights
    soffice_cmd = "soffice"
    if not shutil.which("soffice"):
        win_path = r"C:\Program Files\LibreOffice\program\soffice.exe"
        if Path(win_path).exists():
            soffice_cmd = win_path

    if not soffice_cmd or not shutil.which(soffice_cmd):
        print("LibreOffice not found. Cannot perform comparison.")
        return

    print(f"Running LibreOffice conversion for {input_pptx}...")
    try:
        subprocess.run([soffice_cmd, "--headless", "--convert-to", "pptx", str(input_pptx), "--outdir", str(out_dir)], check=True)
    except subprocess.CalledProcessError as e:
        print(f"LibreOffice conversion failed: {e}")
        return

    actual_pptx = out_dir / input_pptx.name

    try:
        prs_est = Presentation(str(input_pptx))
        prs_act = Presentation(str(actual_pptx))
    except Exception as e:
        print(f"Error reading presentations: {e}")
        return

    if slide_idx < 1 or slide_idx > len(prs_est.slides):
        print(f"Slide index {slide_idx} out of range.")
        return

    slide_est = prs_est.slides[slide_idx - 1]
    slide_act = prs_act.slides[slide_idx - 1]

    print(f"\n--- Height Comparison on Slide {slide_idx} ---")
    print("| Element Name (Text Preview) | Estimated Height (in) | Actual Height (in) | Difference (in) |")
    print("|:---|---:|---:|---:|")
    
    for shape_est, shape_act in zip(slide_est.shapes, slide_act.shapes):
        if not shape_est.has_text_frame:
            continue
            
        text = shape_est.text[:15].replace("\n", " ").replace("\x0b", " ")
        est_h = shape_est.height.inches if shape_est.height else 0
        act_h = shape_act.height.inches if shape_act.height else 0
        diff = act_h - est_h
        
        print(f"| {text} | {est_h:.3f} | {act_h:.3f} | {diff:+.3f} |")

    print(f"\nComparison complete. Checked against {actual_pptx}")
