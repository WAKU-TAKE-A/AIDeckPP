from pptx import Presentation
prs = Presentation('Inputs/template_tate_01.pptx')
slide = prs.slides[0]
print('Shape analysis on Slide 0:')
for i, shape in enumerate(slide.shapes):
    if shape.has_text_frame:
        lines = shape.text.count('\n') + 1
        font_size = None
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                if run.font.size:
                    font_size = run.font.size.pt
                    break
            if font_size: break
        print(f'Shape {i}: lines={lines}, font_size={font_size}pt, text="{shape.text[:20]}..."')
