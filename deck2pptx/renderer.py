from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pathlib import Path
from .models import Deck, Slide, Text, BulletList, Image, Table, Gallery, Flow
from .layout import Layout, get_slide_layout_type
from .theme import Theme

def render_deck(deck: Deck, output_path: str | Path, base_dir: str | Path = '.'):
    prs = Presentation()
    base_dir = Path(base_dir)
    theme = Theme(deck.theme)
    
    # Set slide dimensions
    if deck.orientation == 'portrait':
        prs.slide_width = Inches(7.5)
        prs.slide_height = Inches(13.333)
    else:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
    layout = Layout(prs.slide_width, prs.slide_height)
    
    # Use blank layout (index 6 is blank by default in standard templates)
    slide_layout = prs.slide_layouts[6]
    
    for slide_model in deck.slides:
        slide = prs.slides.add_slide(slide_layout)
        
        # Add slide notes if any
        if slide_model.notes and slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_model.notes
            
        layout_type = get_slide_layout_type(slide_model)
        
        # Render Title/Subtitle
        if layout_type == "title":
            title_y = layout.slide_height / 3
            txBox = slide.shapes.add_textbox(layout.title_x, title_y, layout.title_width, Inches(1.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_model.title
            p.alignment = PP_ALIGN.CENTER
            p.font.name = theme.font_name
            p.font.size = theme.size_title
            p.font.color.rgb = theme.color_text
            
            if slide_model.subtitle:
                p2 = tf.add_paragraph()
                p2.text = slide_model.subtitle
                p2.alignment = PP_ALIGN.CENTER
                p2.font.name = theme.font_name
                p2.font.size = theme.size_subtitle
                p2.font.color.rgb = theme.color_text_light
        else:
            if slide_model.title:
                txBox = slide.shapes.add_textbox(layout.title_x, layout.title_y, layout.title_width, layout.title_height)
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = slide_model.title
                p.font.name = theme.font_name
                p.font.size = theme.size_title
                p.font.color.rgb = theme.color_primary
                
                if slide_model.subtitle:
                    p2 = tf.add_paragraph()
                    p2.text = slide_model.subtitle
                    p2.font.name = theme.font_name
                    p2.font.size = theme.size_subtitle
                    p2.font.color.rgb = theme.color_text_light
        
        # Simple rendering layout flow
        current_y = layout.content_y
        content_x = layout.content_x
        for element in slide_model.elements:
            if isinstance(element, Text):
                txBox = slide.shapes.add_textbox(content_x, current_y, layout.content_width, Inches(1))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = element.content
                p.font.name = theme.font_name
                p.font.size = theme.size_body
                p.font.color.rgb = theme.color_text
                current_y += Inches(0.5)
                
            elif isinstance(element, BulletList):
                txBox = slide.shapes.add_textbox(content_x, current_y, layout.content_width, Inches(2))
                tf = txBox.text_frame
                tf.word_wrap = True
                for i, item in enumerate(element.items):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = item
                    if i > 0:
                        p.level = 1
                    p.font.name = theme.font_name
                    p.font.size = theme.size_body if i == 0 else theme.size_body_small
                    p.font.color.rgb = theme.color_text
                current_y += Inches(2)
                
            elif isinstance(element, Image):
                img_path = base_dir / element.source
                try:
                    # contain fit to width if needed
                    img_width = Inches(4) if Inches(4) < layout.content_width else layout.content_width
                    slide.shapes.add_picture(str(img_path), content_x, current_y, width=img_width)
                    current_y += Inches(3)
                except Exception as e:
                    print(f"Failed to load image {img_path}: {e}")
                    
            elif isinstance(element, Table):
                rows = len(element.rows) + 1 if element.headers else len(element.rows)
                cols = len(element.headers) if element.headers else (len(element.rows[0]) if element.rows else 1)
                
                table_shape = slide.shapes.add_table(rows, cols, content_x, current_y, layout.content_width, Inches(1))
                table = table_shape.table
                
                row_offset = 0
                if element.headers:
                    for i, header in enumerate(element.headers):
                        table.cell(0, i).text = header
                    row_offset = 1
                        
                for r_idx, row in enumerate(element.rows):
                    for c_idx, cell in enumerate(row):
                        cell_obj = table.cell(r_idx + row_offset, c_idx)
                        cell_obj.text = str(cell)
                        for p in cell_obj.text_frame.paragraphs:
                            p.font.name = theme.font_name
                            p.font.size = theme.size_body_small
                        
                current_y += Inches(2)
                
            elif isinstance(element, Gallery):
                num_images = len(element.images)
                if num_images == 0:
                    continue
                elif num_images == 1:
                    cols, rows_ct = 1, 1
                elif num_images == 2:
                    cols, rows_ct = 2, 1
                elif num_images == 3:
                    cols, rows_ct = 3, 1
                elif num_images == 4:
                    cols, rows_ct = 2, 2
                else:
                    cols, rows_ct = 3, 2
                    
                img_width = layout.content_width / cols
                img_height = layout.content_height / rows_ct
                
                for i, img in enumerate(element.images[:cols*rows_ct]):
                    r = i // cols
                    c = i % cols
                    x = content_x + (c * img_width)
                    y = current_y + (r * img_height)
                    img_path = base_dir / img.source
                    try:
                        slide.shapes.add_picture(str(img_path), x, y, width=img_width)
                    except Exception as e:
                        print(f"Failed to load image {img_path}: {e}")
                
                current_y += Inches(4.5)
                
            elif isinstance(element, Flow):
                node_width = Inches(1.5)
                node_height = Inches(0.5)
                
                node_positions = {}
                
                if element.direction == 'horizontal':
                    for i, node in enumerate(element.nodes):
                        x = content_x + i * (node_width + Inches(0.5))
                        y = current_y
                        shape = slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE,
                            x, y, node_width, node_height
                        )
                        shape.text = node.label
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = theme.color_accent
                        shape.line.color.rgb = theme.color_accent
                        for p in shape.text_frame.paragraphs:
                            p.font.name = theme.font_name
                            p.font.size = theme.size_body_small
                        node_positions[node.id] = (x + node_width, y + node_height / 2) # right center for out, left center for in
                        # Actually let's just store the right-center and left-center
                        node_positions[f"{node.id}_out"] = (x + node_width, y + node_height / 2)
                        node_positions[f"{node.id}_in"] = (x, y + node_height / 2)
                        
                    for edge in element.edges:
                        if f"{edge.from_node}_out" in node_positions and f"{edge.to_node}_in" in node_positions:
                            fx, fy = node_positions[f"{edge.from_node}_out"]
                            tx, ty = node_positions[f"{edge.to_node}_in"]
                            
                            if tx >= fx:
                                arrow_x = fx
                                arrow_width = max(tx - fx, Inches(0.1))
                                shape_type = MSO_SHAPE.RIGHT_ARROW
                            else:
                                arrow_x = tx
                                arrow_width = max(fx - tx, Inches(0.1))
                                shape_type = MSO_SHAPE.LEFT_ARROW
                                
                            slide.shapes.add_shape(
                                shape_type,
                                arrow_x, fy - Inches(0.1), arrow_width, Inches(0.2)
                            )
                else: # vertical
                    for i, node in enumerate(element.nodes):
                        x = content_x
                        y = current_y + i * (node_height + Inches(0.5))
                        shape = slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE,
                            x, y, node_width, node_height
                        )
                        shape.text = node.label
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = theme.color_accent
                        shape.line.color.rgb = theme.color_accent
                        for p in shape.text_frame.paragraphs:
                            p.font.name = theme.font_name
                            p.font.size = theme.size_body_small
                        node_positions[f"{node.id}_out"] = (x + node_width / 2, y + node_height)
                        node_positions[f"{node.id}_in"] = (x + node_width / 2, y)
                
                    for edge in element.edges:
                        if f"{edge.from_node}_out" in node_positions and f"{edge.to_node}_in" in node_positions:
                            fx, fy = node_positions[f"{edge.from_node}_out"]
                            tx, ty = node_positions[f"{edge.to_node}_in"]
                            # Draw down/up arrow
                            if ty >= fy:
                                arrow_y = fy
                                arrow_height = max(ty - fy, Inches(0.1))
                                shape_type = MSO_SHAPE.DOWN_ARROW
                            else:
                                arrow_y = ty
                                arrow_height = max(fy - ty, Inches(0.1))
                                shape_type = MSO_SHAPE.UP_ARROW
                                
                            slide.shapes.add_shape(
                                shape_type,
                                fx - Inches(0.1), arrow_y, Inches(0.2), arrow_height
                            )

                current_y += Inches(3)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
