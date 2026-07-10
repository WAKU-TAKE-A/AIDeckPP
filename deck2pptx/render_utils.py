from pptx.util import Pt
from .text_utils import parse_inline_formatting

def _set_text_frame_text(text_frame, content, font_name=None, font_size=None, font_color=None, input_format="markdown"):
    text_frame.clear()
    lines = str(content).split('\n')
    for i, line in enumerate(lines):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        
        chunks = parse_inline_formatting(line, input_format=input_format)
        for chunk in chunks:
            run = p.add_run()
            run.text = chunk.text
            if chunk.bold:
                run.font.bold = True
            if chunk.italic:
                run.font.italic = True
            if font_name:
                run.font.name = font_name
            if font_size:
                run.font.size = font_size
            if font_color:
                run.font.color.rgb = font_color
                
        p.space_before = Pt(0)
        p.space_after = Pt(0)

