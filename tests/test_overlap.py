import pytest
from pptx.util import Inches
from deck2pptx.models import Deck, Slide, Text, Image, BulletList
from deck2pptx.renderer import render_deck
from pptx import Presentation

def test_no_overlap_text_image_text(tmp_path):
    # Dummy image
    img_path = tmp_path / "dummy.jpg"
    from PIL import Image as PILImage
    PILImage.new('RGB', (100, 100)).save(img_path)

    deck = Deck(
        slides=[
            Slide(
                title="Test Slide",
                elements=[
                    Text(content="First text block\nLine 2\nLine 3"),
                    Image(source=str(img_path)),
                    Text(content="Second text block")
                ]
            )
        ]
    )

    out_path = tmp_path / "out.pptx"
    render_deck(deck, str(out_path), base_dir=tmp_path)

    prs = Presentation(str(out_path))
    slide = prs.slides[0]
    
    content_shapes = [s for s in slide.shapes if not s.is_placeholder and not (s.has_text_frame and s.text.startswith("Test Slide"))]
    assert len(content_shapes) == 3

    for i in range(len(content_shapes) - 1):
        s1 = content_shapes[i]
        s2 = content_shapes[i + 1]
        assert s2.top >= s1.top + s1.height

def test_no_overlap_bullet_text(tmp_path):
    deck = Deck(
        slides=[
            Slide(
                title="Test Slide",
                elements=[
                    BulletList(items=["Item 1", "Item 2", "Item 3"]),
                    Text(content="Follow up text")
                ]
            )
        ]
    )

    out_path = tmp_path / "out2.pptx"
    render_deck(deck, str(out_path), base_dir=tmp_path)

    prs = Presentation(str(out_path))
    slide = prs.slides[0]
    
    content_shapes = [s for s in slide.shapes if not s.is_placeholder and not (s.has_text_frame and s.text.startswith("Test Slide"))]
    assert len(content_shapes) == 2

    s1 = content_shapes[0]
    s2 = content_shapes[1]
    assert s2.top >= s1.top + s1.height

def test_height_hint_override(tmp_path):
    deck = Deck(
        slides=[
            Slide(
                title="Test Hint",
                elements=[
                    Text(content="First text block", height_hint=Inches(3)),
                    Text(content="Second text block")
                ]
            )
        ]
    )

    out_path = tmp_path / "out3.pptx"
    render_deck(deck, str(out_path), base_dir=tmp_path)

    prs = Presentation(str(out_path))
    slide = prs.slides[0]
    
    content_shapes = [s for s in slide.shapes if not s.is_placeholder and not (s.has_text_frame and s.text.startswith("Test Hint"))]
    assert len(content_shapes) == 2

    s1 = content_shapes[0]
    s2 = content_shapes[1]
    assert s1.height == Inches(3)
    assert s2.top >= s1.top + s1.height
